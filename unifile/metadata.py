"""UniFile — Metadata extraction from files (EXIF, audio, documents, archives)."""
import os, re, json, hashlib, base64, io, zipfile, gzip, shutil, subprocess
from collections import Counter
from pathlib import Path
import xml.etree.ElementTree as ET

# Extension sets — kept here instead of files.py so metadata.py stays
# self-contained even when files.py hasn't been imported yet.
_META_IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff',
                    '.webp', '.heic', '.heif', '.avif', '.jxl',
                    '.cr2', '.cr3', '.nef', '.arw', '.dng', '.orf', '.rw2',
                    '.pef', '.srw', '.raw'}
_META_AUDIO_EXTS = {'.mp3', '.wav', '.flac', '.aac', '.ogg', '.wma', '.m4a',
                    '.opus', '.aiff', '.aif', '.ape', '.mka', '.wv', '.tta',
                    '.dsf', '.dff', '.caf', '.mid', '.midi'}
_META_VIDEO_EXTS = {'.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv', '.webm',
                    '.m4v', '.mpg', '.mpeg', '.3gp', '.ts', '.mts', '.m2ts',
                    '.vob', '.ogv', '.asf', '.f4v', '.h264', '.h265', '.hevc'}
_META_PDF_EXTS = {'.pdf'}
_META_DOCX_EXTS = {'.docx'}
_META_XLSX_EXTS = {'.xlsx'}
_META_PPTX_EXTS = {'.pptx'}

from unifile.bootstrap import (
    HAS_PILLOW, HAS_PILLOW_HEIF, HAS_EXIFREAD, HAS_MUTAGEN,
    HAS_PYPDF, HAS_PYTHON_DOCX, HAS_OPENPYXL, HAS_PYTHON_PPTX,
    HAS_PSD_TOOLS, HAS_RARFILE, HAS_PY7ZR
)
from unifile.config import _APP_DATA_DIR
try:
    from PIL import Image as _PILImage
    from PIL.ExifTags import TAGS as _EXIF_TAGS, GPSTAGS as _GPS_TAGS
except ImportError:
    pass
try:
    import exifread as _exifread
except ImportError:
    pass
try:
    import mutagen as _mutagen
    from mutagen.easyid3 import EasyID3 as _EasyID3
    from mutagen.mp3 import MP3 as _MP3
    from mutagen.flac import FLAC as _FLAC
    from mutagen.mp4 import MP4 as _MP4
    from mutagen.oggvorbis import OggVorbis as _OggVorbis
except ImportError:
    pass
try:
    from pypdf import PdfReader as _PdfReader
except ImportError:
    pass
try:
    from docx import Document as _DocxDocument
except ImportError:
    pass
try:
    from openpyxl import load_workbook as _load_workbook
except ImportError:
    pass
try:
    from pptx import Presentation as _PptxPresentation
except ImportError:
    pass
try:
    import psd_tools as _psd_tools
except ImportError:
    pass
try:
    import rarfile as _rarfile
except ImportError:
    pass
try:
    import py7zr as _py7zr
except ImportError:
    pass

# ── Level 3: Metadata extraction ─────────────────────────────────────────────

def extract_prproj_metadata(filepath: str) -> list:
    """Extract sequence/clip names from .prproj files (gzipped XML).
    Returns list of name strings useful for categorization."""
    names = []
    try:
        with gzip.open(filepath, 'rt', encoding='utf-8', errors='ignore') as f:
            content = f.read(500_000)  # Cap read to 500KB

        # Search for descriptive elements in Premiere XML
        # Match content inside common name/title tags
        for pattern in [
            r'<(?:Name|Title|SequenceName|ActualName|n)>(.*?)</(?:Name|Title|SequenceName|ActualName|n)>',
            r'ObjectURef="([^"]+)"',
            r'<Label>(.*?)</Label>',
        ]:
            for match in re.finditer(pattern, content, re.IGNORECASE):
                name = match.group(1).strip()
                if (name and len(name) > 3
                    and not name.startswith(('Sequence', 'Untitled', 'Comp'))
                    and not name.isdigit()
                    and not re.match(r'^[0-9a-f\-]{20,}$', name)):  # Skip UUIDs
                    names.append(name)
    except Exception:
        pass
    return names[:20]  # Cap results


def extract_psd_metadata(filepath: str) -> list:
    """Extract layer names and metadata from .psd files via psd-tools.
    Falls back gracefully if psd-tools not installed."""
    if not HAS_PSD_TOOLS:
        return []
    names = []
    try:
        psd = _psd_tools.PSDImage.open(filepath)
        for layer in psd.descendants():
            if hasattr(layer, 'name') and layer.name:
                name = layer.name.strip()
                # Skip generic Photoshop layer names
                if name.lower() not in {'layer 1', 'layer 2', 'layer 3', 'background',
                                          'group 1', 'group 2', 'copy', 'shape 1', 'shape 2',
                                          'layer', 'group', 'effect', 'mask'}:
                    names.append(name)
        # Also extract document info if available
        if hasattr(psd, 'image_resources'):
            pass  # Could extract title from IPTC/XMP but keeping it simple
    except Exception:
        pass
    return names[:30]  # Cap results


# Envato item code pattern: 7-8 digit numbers common in Envato downloads
_ENVATO_ID_PATTERN = re.compile(r'(?:^|[\-_\s])(\d{7,8})(?:[\-_\s]|$)')

def detect_envato_item_code(folder_name: str) -> str:
    """Detect Envato marketplace item codes in folder names.
    Returns the item code string or empty string."""
    match = _ENVATO_ID_PATTERN.search(folder_name)
    return match.group(1) if match else ''


def extract_folder_metadata(folder_path: str, log_cb=None) -> dict:
    """Extract all available metadata from files inside a folder.
    Returns dict with keywords, project_names, envato_id, etc."""
    metadata = {
        'keywords': [],         # Keywords extracted from file metadata
        'project_names': [],    # Named sequences, compositions, etc.
        'envato_id': '',        # Envato item code if detected
        'primary_app': '',      # Detected primary Adobe app
        'has_aep': False,
        'has_prproj': False,
        'has_psd': False,
        'has_mogrt': False,
    }

    folder_p = Path(folder_path)

    # Detect Envato item code from folder name
    metadata['envato_id'] = detect_envato_item_code(folder_p.name)

    # Scan files for metadata (limit depth and count for performance)
    scanned = 0
    max_scan = 10  # Max files to parse metadata from

    try:
        for item in folder_p.rglob('*'):
            if not item.is_file():
                continue
            ext = item.suffix.lower()

            # Track what app-specific files exist
            if ext in ('.aep', '.aet'):
                metadata['has_aep'] = True
                metadata['primary_app'] = metadata['primary_app'] or 'After Effects'
            elif ext == '.prproj':
                metadata['has_prproj'] = True
                metadata['primary_app'] = metadata['primary_app'] or 'Premiere Pro'
                if scanned < max_scan:
                    names = extract_prproj_metadata(str(item))
                    metadata['project_names'].extend(names)
                    scanned += 1
            elif ext in ('.psd', '.psb'):
                metadata['has_psd'] = True
                metadata['primary_app'] = metadata['primary_app'] or 'Photoshop'
                if scanned < max_scan and HAS_PSD_TOOLS:
                    names = extract_psd_metadata(str(item))
                    metadata['keywords'].extend(names)
                    scanned += 1
            elif ext == '.mogrt':
                metadata['has_mogrt'] = True
                metadata['primary_app'] = metadata['primary_app'] or 'After Effects'

            if scanned >= max_scan:
                break
    except (PermissionError, OSError):
        pass

    return metadata



# ── Level 3.5: Envato API metadata enrichment ────────────────────────────────

# Envato API category → UniFile category mapping
_ENVATO_CAT_MAP = {
    'after-effects-project-files': 'After Effects - Templates',
    'after-effects-presets': 'After Effects - Presets & Scripts',
    'after-effects-scripts': 'After Effects - Presets & Scripts',
    'premiere-pro-templates': 'Premiere Pro - Templates',
    'premiere-pro-presets': 'Premiere Pro - Presets & Effects',
    'motion-graphics': 'Motion Graphics',
    'stock-footage': 'Stock Footage - General',
    'stock-music': 'Stock Music & Audio',
    'sound-effects': 'Sound Effects & SFX',
    'fonts': 'Fonts & Typography',
    'graphics': 'Illustrator - Vectors & Assets',
    'add-ons': 'Photoshop - Actions',
    'photos': 'Stock Photos - General',
    'video-templates': 'Video Editing - General',
    'presentation-templates': 'Presentations & PowerPoint',
    '3d': '3D - Models & Objects',
    'logos': 'Logo & Identity',
    'product-mockups': 'Photoshop - Mockups',
    'infographics': 'Infographic',
    'web-templates': 'UI & UX Design',
    'backgrounds': 'Backgrounds & Textures',
    'textures': 'Photoshop - Patterns',
    'icons': 'Illustrator - Icons & UI Kits',
    'illustrations': 'Clipart & Illustrations',
}

# Persistent API key storage path
_ENVATO_KEY_FILE = os.path.join(_APP_DATA_DIR, 'envato_api_key.txt')

def _load_envato_api_key() -> str:
    """Load Envato API key from file. Returns empty string if not set."""
    try:
        with open(_ENVATO_KEY_FILE, 'r') as f:
            return f.read().strip()
    except (FileNotFoundError, OSError):
        return ''

def _save_envato_api_key(key: str):
    """Save Envato API key to file."""
    try:
        with open(_ENVATO_KEY_FILE, 'w') as f:
            f.write(key.strip())
    except OSError:
        pass

# Simple in-memory cache for API responses
_envato_cache = {}

def _envato_api_classify(item_id: str) -> tuple:
    """Query Envato Market API to get item category and tags.
    Returns (category, confidence, detail) or (None, 0, '').
    Requires API key stored in envato_api_key.txt alongside the script."""
    from unifile.classifier import categorize_folder
    if not item_id:
        return (None, 0, '')

    api_key = _load_envato_api_key()
    if not api_key:
        return (None, 0, '')

    # Check cache
    if item_id in _envato_cache:
        return _envato_cache[item_id]

    try:
        import urllib.request, urllib.error
        url = f"https://api.envato.com/v3/market/catalog/item?id={item_id}"
        req = urllib.request.Request(url, headers={
            'Authorization': f'Bearer {api_key}',
            'User-Agent': 'UniFile/3.0'
        })
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode('utf-8'))

        # Extract classification signals
        api_cat = data.get('classification', '')
        api_tags = data.get('tags', [])
        api_name = data.get('name', '')

        # Try direct category mapping
        cat_slug = api_cat.lower().replace(' ', '-') if api_cat else ''
        for envato_key, fo_cat in _ENVATO_CAT_MAP.items():
            if envato_key in cat_slug:
                result = (fo_cat, 88, f"envato_api:{item_id}→{api_cat}")
                _envato_cache[item_id] = result
                return result

        # Fallback: classify the API-provided item name through keyword engine
        if api_name:
            kw_cat, kw_conf, _ = categorize_folder(api_name)
            if kw_cat and kw_conf >= 40:
                result = (kw_cat, min(kw_conf + 15, 92), f"envato_name:\"{api_name}\"")
                _envato_cache[item_id] = result
                return result

        # Fallback: try classifying tags
        for tag in api_tags[:5]:
            t_cat, t_conf, _ = categorize_folder(tag)
            if t_cat and t_conf >= 60:
                result = (t_cat, min(t_conf + 5, 85), f"envato_tag:\"{tag}\"")
                _envato_cache[item_id] = result
                return result

    except Exception:
        pass

    _envato_cache[item_id] = (None, 0, '')
    return (None, 0, '')



# ── Content extraction for intelligent file renaming ─────────────────────────
_CONTENT_TEXT_EXTS = {
    '.py', '.ps1', '.psm1', '.psd1', '.js', '.ts', '.jsx', '.tsx', '.c', '.cpp',
    '.h', '.hpp', '.cs', '.java', '.go', '.rs', '.rb', '.php', '.swift', '.kt',
    '.sh', '.bash', '.zsh', '.bat', '.cmd', '.vbs',
    '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.env',
    '.xml', '.svg', '.sql', '.graphql', '.proto',
    '.txt', '.md', '.rst', '.log', '.nfo', '.readme',
    '.csv', '.tsv', '.tex', '.r', '.m', '.lua', '.pl', '.pm',
    '.dockerfile', '.makefile', '.cmake',
}
_CONTENT_AUDIO_EXTS = {
    '.mp3', '.flac', '.ogg', '.m4a', '.wma', '.aac', '.opus', '.wav', '.aiff',
}

_CONTENT_MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB safety limit


def _extract_file_content(file_path: str, max_chars: int = 800) -> str:
    """Extract text content from a file for LLM-based intelligent naming.

    Returns a sanitized text snippet (up to max_chars), or empty string on failure.
    Never raises exceptions. Supports PDF, Word, Excel, PowerPoint, PSD, audio, HTML,
    and plain text/code files.
    """
    try:
        fsize = os.path.getsize(file_path)
        if fsize == 0 or fsize > _CONTENT_MAX_FILE_SIZE:
            return ''
    except OSError:
        return ''

    ext = os.path.splitext(file_path)[1].lower()
    content = ''

    # ── PDF ───────────────────────────────────────────────────────────────
    if ext == '.pdf' and HAS_PYPDF:
        try:
            reader = _PdfReader(file_path)
            parts = []
            for page in reader.pages[:2]:
                txt = page.extract_text()
                if txt:
                    parts.append(txt.strip())
            content = '\n'.join(parts)
        except Exception:
            pass

    # ── Word (.docx) ─────────────────────────────────────────────────────
    elif ext == '.docx' and HAS_PYTHON_DOCX:
        try:
            doc = _DocxDocument(file_path)
            parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            content = '\n'.join(parts)
        except Exception:
            pass

    # ── Excel (.xlsx) ────────────────────────────────────────────────────
    elif ext == '.xlsx' and HAS_OPENPYXL:
        try:
            wb = _load_workbook(file_path, read_only=True, data_only=True)
            parts = [f"Sheets: {', '.join(wb.sheetnames)}"]
            ws = wb.worksheets[0]
            for row in ws.iter_rows(max_row=5, values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                parts.append(' | '.join(cells))
            wb.close()
            content = '\n'.join(parts)
        except Exception:
            pass

    # ── PowerPoint (.pptx) ───────────────────────────────────────────────
    elif ext == '.pptx' and HAS_PYTHON_PPTX:
        try:
            prs = _PptxPresentation(file_path)
            parts = []
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            parts.append(txt)
            content = '\n'.join(parts)
        except Exception:
            pass

    # ── Photoshop (.psd) ─────────────────────────────────────────────────
    elif ext == '.psd' and HAS_PSD_TOOLS:
        try:
            psd = _psd_tools.PSDImage.open(file_path)
            layer_names = [layer.name for layer in psd.descendants() if layer.name]
            content = f"Layers: {', '.join(layer_names[:20])}"
        except Exception:
            pass

    # ── Audio files (metadata tags) ──────────────────────────────────────
    elif ext in _CONTENT_AUDIO_EXTS and HAS_MUTAGEN:
        try:
            audio = _mutagen.File(file_path, easy=True)
            if audio:
                parts = []
                for tag in ('artist', 'title', 'album', 'genre'):
                    val = audio.get(tag)
                    if val:
                        parts.append(f"{tag}: {val[0] if isinstance(val, list) else val}")
                content = ' | '.join(parts)
        except Exception:
            pass

    # ── HTML files ────────────────────────────────────────────────────────
    elif ext in ('.html', '.htm'):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                raw = f.read(max_chars * 3)
            # Extract title
            title_m = re.search(r'<title[^>]*>(.*?)</title>', raw, re.IGNORECASE | re.DOTALL)
            title = title_m.group(1).strip() if title_m else ''
            # Strip tags for body text
            body = re.sub(r'<[^>]+>', ' ', raw)
            body = re.sub(r'\s+', ' ', body).strip()
            content = f"Title: {title}\n{body}" if title else body
        except Exception:
            pass

    # ── Plain text / code / config ───────────────────────────────────────
    elif ext in _CONTENT_TEXT_EXTS or ext == '':
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(max_chars + 200)
        except Exception:
            pass

    if not content or not content.strip():
        return ''

    # Sanitize: strip LLM injection characters and truncate
    content = re.sub(r'[{}\[\]<>]', '', content)
    content = re.sub(r'\s+', ' ', content).strip()
    return content[:max_chars]



class MetadataExtractor:
    """Extracts file metadata for PC File Organizer (Mode 3).

    Returns a dict with standardised keys. Missing fields are omitted.
    All operations are wrapped in try/except — never crashes the caller.
    """

    @staticmethod
    def capabilities() -> dict:
        """Return which extractors are available based on installed libraries."""
        return {
            'images':  HAS_PILLOW or HAS_EXIFREAD,
            'audio':   HAS_MUTAGEN,
            'video':   shutil.which('ffprobe') is not None,
            'pdf':     HAS_PYPDF,
            'docx':    HAS_PYTHON_DOCX,
            'xlsx':    HAS_OPENPYXL,
        }

    @staticmethod
    def extract(filepath: str, log_cb=None) -> dict:
        """Extract metadata from a file. Returns dict (may be empty)."""
        meta = {}
        try:
            ext = os.path.splitext(filepath)[1].lower()
            if ext in _META_IMAGE_EXTS:
                meta = MetadataExtractor._extract_image(filepath)
            elif ext in _META_AUDIO_EXTS:
                meta = MetadataExtractor._extract_audio(filepath)
            elif ext in _META_VIDEO_EXTS:
                meta = MetadataExtractor._extract_video(filepath)
            elif ext in _META_PDF_EXTS:
                meta = MetadataExtractor._extract_pdf(filepath)
            elif ext in _META_DOCX_EXTS:
                meta = MetadataExtractor._extract_docx(filepath)
            elif ext in _META_XLSX_EXTS:
                meta = MetadataExtractor._extract_xlsx(filepath)
            elif ext in _META_PPTX_EXTS:
                meta = MetadataExtractor._extract_pptx(filepath)
        except Exception as e:
            if log_cb:
                log_cb(f"    [META] Error extracting {os.path.basename(filepath)}: {e}")
        return meta

    @staticmethod
    def format_summary(meta: dict) -> str:
        """Build a human-readable one-line summary from metadata dict."""
        if not meta:
            return ""
        parts = []
        mt = meta.get('_type', '')
        if mt == 'image':
            if meta.get('width') and meta.get('height'):
                parts.append(f"{meta['width']}×{meta['height']}")
            if meta.get('camera_make') or meta.get('camera_model'):
                cam = ' '.join(filter(None, [meta.get('camera_make', ''),
                                              meta.get('camera_model', '')]))
                parts.append(cam.strip())
            if meta.get('date_taken'):
                parts.append(meta['date_taken'][:10])
            if meta.get('gps_lat') is not None and meta.get('gps_lon') is not None:
                parts.append(f"GPS:{meta['gps_lat']:.4f},{meta['gps_lon']:.4f}")
        elif mt == 'audio':
            if meta.get('artist'):
                parts.append(meta['artist'])
            if meta.get('album'):
                parts.append(meta['album'])
            if meta.get('title'):
                parts.append(f"\"{meta['title']}\"")
            if meta.get('duration'):
                dur = int(meta['duration'])
                parts.append(f"{dur // 60}:{dur % 60:02d}")
            if meta.get('bitrate'):
                parts.append(f"{meta['bitrate']}kbps")
        elif mt == 'video':
            if meta.get('width') and meta.get('height'):
                parts.append(f"{meta['width']}×{meta['height']}")
            if meta.get('codec'):
                parts.append(meta['codec'])
            if meta.get('duration'):
                dur = int(meta['duration'])
                m, s = divmod(dur, 60)
                h, m = divmod(m, 60)
                parts.append(f"{h}:{m:02d}:{s:02d}" if h else f"{m}:{s:02d}")
            if meta.get('fps'):
                parts.append(f"{meta['fps']:.1f}fps")
        elif mt == 'pdf':
            if meta.get('title'):
                parts.append(f"\"{meta['title'][:40]}\"")
            if meta.get('author'):
                parts.append(f"by {meta['author']}")
            if meta.get('pages'):
                parts.append(f"{meta['pages']} pg")
        elif mt in ('docx', 'xlsx', 'pptx'):
            if meta.get('title'):
                parts.append(f"\"{meta['title'][:40]}\"")
            if meta.get('author'):
                parts.append(f"by {meta['author']}")
            if meta.get('last_modified_by'):
                parts.append(f"mod: {meta['last_modified_by']}")
        return ' | '.join(parts) if parts else ''

    @staticmethod
    def format_tooltip(meta: dict) -> str:
        """Build a detailed multi-line tooltip from metadata dict."""
        if not meta:
            return ""
        lines = []
        mt = meta.get('_type', '')
        skip = {'_type'}
        _LABELS = {
            'width': 'Width', 'height': 'Height', 'camera_make': 'Camera Make',
            'camera_model': 'Camera Model', 'date_taken': 'Date Taken',
            'gps_lat': 'Latitude', 'gps_lon': 'Longitude', 'gps_alt': 'Altitude',
            'orientation': 'Orientation', 'color_space': 'Color Space',
            'artist': 'Artist', 'album': 'Album', 'title': 'Title',
            'year': 'Year', 'genre': 'Genre', 'track': 'Track',
            'duration': 'Duration', 'bitrate': 'Bitrate', 'sample_rate': 'Sample Rate',
            'channels': 'Channels', 'codec': 'Codec', 'fps': 'FPS',
            'video_bitrate': 'Video Bitrate', 'audio_codec': 'Audio Codec',
            'author': 'Author', 'subject': 'Subject', 'creator': 'Creator',
            'pages': 'Pages', 'creation_date': 'Created',
            'last_modified_by': 'Last Modified By', 'modified': 'Modified',
            'created': 'Created', 'revision': 'Revision', 'keywords': 'Keywords',
            'sheet_count': 'Sheets', 'slide_count': 'Slides',
        }
        for k, v in meta.items():
            if k in skip or v is None or v == '':
                continue
            label = _LABELS.get(k, k.replace('_', ' ').title())
            if k == 'duration' and isinstance(v, (int, float)):
                dur = int(v)
                m, s = divmod(dur, 60)
                h, m2 = divmod(m, 60)
                v = f"{h}:{m2:02d}:{s:02d}" if h else f"{m}:{s:02d}"
            elif k == 'bitrate' and isinstance(v, (int, float)):
                v = f"{int(v)} kbps"
            elif k == 'sample_rate' and isinstance(v, (int, float)):
                v = f"{int(v)} Hz"
            elif k in ('gps_lat', 'gps_lon') and isinstance(v, float):
                v = f"{v:.6f}"
            lines.append(f"{label}: {v}")
        return '\n'.join(lines)

    # ── Image metadata ────────────────────────────────────────────────────────
    @staticmethod
    def _extract_image(filepath: str) -> dict:
        meta = {'_type': 'image'}
        # Try Pillow first (most common formats)
        if HAS_PILLOW:
            try:
                with _PILImage.open(filepath) as img:
                    meta['width'] = img.width
                    meta['height'] = img.height
                    # EXIF data
                    exif_data = img.getexif()
                    if exif_data:
                        for tag_id, val in exif_data.items():
                            tag_name = _EXIF_TAGS.get(tag_id, str(tag_id))
                            if tag_name == 'Make':
                                meta['camera_make'] = str(val).strip()
                            elif tag_name == 'Model':
                                meta['camera_model'] = str(val).strip()
                            elif tag_name == 'DateTimeOriginal':
                                meta['date_taken'] = str(val).strip()
                            elif tag_name == 'DateTime' and 'date_taken' not in meta:
                                meta['date_taken'] = str(val).strip()
                            elif tag_name == 'Orientation':
                                meta['orientation'] = int(val) if val else None
                            elif tag_name == 'ColorSpace':
                                meta['color_space'] = str(val).strip()
                        # GPS info (IFD 0x8825)
                        gps_ifd = exif_data.get_ifd(0x8825)
                        if gps_ifd:
                            meta.update(MetadataExtractor._parse_gps(gps_ifd))
            except Exception:
                pass
        # Fallback to exifread for RAW/HEIC formats (skip for formats Pillow handles well)
        _PILLOW_NATIVE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp'}
        ext_lower = os.path.splitext(filepath)[1].lower()
        if (HAS_EXIFREAD and ext_lower not in _PILLOW_NATIVE_EXTS
                and ('date_taken' not in meta or 'width' not in meta)):
            try:
                with open(filepath, 'rb') as f:
                    tags = _exifread.process_file(f, stop_tag='UNDEF', details=False)
                if tags:
                    if 'EXIF DateTimeOriginal' in tags and 'date_taken' not in meta:
                        meta['date_taken'] = str(tags['EXIF DateTimeOriginal'])
                    if 'Image Make' in tags and 'camera_make' not in meta:
                        meta['camera_make'] = str(tags['Image Make']).strip()
                    if 'Image Model' in tags and 'camera_model' not in meta:
                        meta['camera_model'] = str(tags['Image Model']).strip()
                    if 'EXIF ExifImageWidth' in tags and 'width' not in meta:
                        try: meta['width'] = int(str(tags['EXIF ExifImageWidth']))
                        except ValueError: pass
                    if 'EXIF ExifImageLength' in tags and 'height' not in meta:
                        try: meta['height'] = int(str(tags['EXIF ExifImageLength']))
                        except ValueError: pass
                    # GPS from exifread
                    if 'GPS GPSLatitude' in tags and 'gps_lat' not in meta:
                        lat = MetadataExtractor._exifread_gps_coord(
                            tags.get('GPS GPSLatitude'), tags.get('GPS GPSLatitudeRef'))
                        lon = MetadataExtractor._exifread_gps_coord(
                            tags.get('GPS GPSLongitude'), tags.get('GPS GPSLongitudeRef'))
                        if lat is not None:
                            meta['gps_lat'] = lat
                        if lon is not None:
                            meta['gps_lon'] = lon
            except Exception:
                pass
        return meta

    @staticmethod
    def _parse_gps(gps_ifd: dict) -> dict:
        """Parse GPS IFD from Pillow EXIF into decimal lat/lon."""
        result = {}
        try:
            lat_data = gps_ifd.get(2)   # GPSLatitude
            lat_ref  = gps_ifd.get(1)   # GPSLatitudeRef ('N' or 'S')
            lon_data = gps_ifd.get(4)   # GPSLongitude
            lon_ref  = gps_ifd.get(3)   # GPSLongitudeRef ('E' or 'W')
            if lat_data and lon_data:
                lat = MetadataExtractor._dms_to_decimal(lat_data, lat_ref)
                lon = MetadataExtractor._dms_to_decimal(lon_data, lon_ref)
                if lat is not None:
                    result['gps_lat'] = lat
                if lon is not None:
                    result['gps_lon'] = lon
            alt_data = gps_ifd.get(6)   # GPSAltitude
            if alt_data is not None:
                try:
                    alt = float(alt_data)
                    alt_ref = gps_ifd.get(5, 0)
                    if alt_ref == 1:
                        alt = -alt
                    result['gps_alt'] = round(alt, 1)
                except (TypeError, ValueError):
                    pass
        except Exception:
            pass
        return result

    @staticmethod
    def _dms_to_decimal(dms_tuple, ref) -> float:
        """Convert (degrees, minutes, seconds) to decimal degrees."""
        try:
            d = float(dms_tuple[0])
            m = float(dms_tuple[1])
            s = float(dms_tuple[2])
            decimal = d + m / 60.0 + s / 3600.0
            if ref in ('S', 'W'):
                decimal = -decimal
            return round(decimal, 6)
        except (TypeError, ValueError, IndexError):
            return None

    @staticmethod
    def _exifread_gps_coord(coord_tag, ref_tag) -> float:
        """Convert exifread GPS coordinate to decimal degrees."""
        try:
            if coord_tag is None:
                return None
            vals = coord_tag.values
            d = float(vals[0].num) / float(vals[0].den) if hasattr(vals[0], 'num') else float(vals[0])
            m = float(vals[1].num) / float(vals[1].den) if hasattr(vals[1], 'num') else float(vals[1])
            s = float(vals[2].num) / float(vals[2].den) if hasattr(vals[2], 'num') else float(vals[2])
            decimal = d + m / 60.0 + s / 3600.0
            if ref_tag and str(ref_tag) in ('S', 'W'):
                decimal = -decimal
            return round(decimal, 6)
        except Exception:
            return None

    # ── Audio metadata ────────────────────────────────────────────────────────
    @staticmethod
    def _extract_audio(filepath: str) -> dict:
        meta = {'_type': 'audio'}
        if not HAS_MUTAGEN:
            return meta
        try:
            mf = _mutagen.File(filepath, easy=True)
            if mf is None:
                return meta
            # Common tags via Easy interface
            for tag, key in [('title', 'title'), ('artist', 'artist'),
                             ('album', 'album'), ('genre', 'genre'),
                             ('date', 'year'), ('tracknumber', 'track')]:
                val = mf.get(tag)
                if val:
                    meta[key] = str(val[0]) if isinstance(val, list) else str(val)
            # Duration and bitrate from info object
            if hasattr(mf, 'info') and mf.info:
                info = mf.info
                if hasattr(info, 'length') and info.length:
                    meta['duration'] = round(info.length, 1)
                if hasattr(info, 'bitrate') and info.bitrate:
                    meta['bitrate'] = int(info.bitrate / 1000)  # kbps
                if hasattr(info, 'sample_rate') and info.sample_rate:
                    meta['sample_rate'] = int(info.sample_rate)
                if hasattr(info, 'channels') and info.channels:
                    meta['channels'] = info.channels
        except Exception:
            # Retry with raw mutagen (non-easy mode) for formats that need it
            try:
                mf = _mutagen.File(filepath)
                if mf and hasattr(mf, 'info') and mf.info:
                    if hasattr(mf.info, 'length') and mf.info.length:
                        meta['duration'] = round(mf.info.length, 1)
                    if hasattr(mf.info, 'bitrate') and mf.info.bitrate:
                        meta['bitrate'] = int(mf.info.bitrate / 1000)
            except Exception:
                pass
        return meta

    # ── Video metadata ────────────────────────────────────────────────────────
    @staticmethod
    def _extract_video(filepath: str) -> dict:
        meta = {'_type': 'video'}
        ffprobe_path = shutil.which('ffprobe')
        if not ffprobe_path:
            return meta
        try:
            cmd = [
                ffprobe_path, '-v', 'quiet',
                '-print_format', 'json',
                '-show_format', '-show_streams',
                filepath
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                return meta
            data = json.loads(result.stdout)

            # Format-level info
            fmt = data.get('format', {})
            if fmt.get('duration'):
                meta['duration'] = round(float(fmt['duration']), 1)

            # Find video and audio streams
            for stream in data.get('streams', []):
                codec_type = stream.get('codec_type', '')
                if codec_type == 'video':
                    if stream.get('width'):
                        meta['width'] = int(stream['width'])
                    if stream.get('height'):
                        meta['height'] = int(stream['height'])
                    if stream.get('codec_name'):
                        meta['codec'] = stream['codec_name']
                    # FPS from r_frame_rate (e.g. "30000/1001")
                    fps_str = stream.get('r_frame_rate', '')
                    if fps_str and '/' in fps_str:
                        num, den = fps_str.split('/')
                        try:
                            fps = float(num) / float(den)
                            if 0 < fps < 1000:
                                meta['fps'] = round(fps, 2)
                        except (ValueError, ZeroDivisionError):
                            pass
                    if stream.get('bit_rate'):
                        meta['video_bitrate'] = int(int(stream['bit_rate']) / 1000)
                elif codec_type == 'audio' and 'audio_codec' not in meta:
                    if stream.get('codec_name'):
                        meta['audio_codec'] = stream['codec_name']
        except (subprocess.TimeoutExpired, json.JSONDecodeError, OSError):
            pass
        return meta

    # ── PDF metadata ──────────────────────────────────────────────────────────
    @staticmethod
    def _extract_pdf(filepath: str) -> dict:
        meta = {'_type': 'pdf'}
        if not HAS_PYPDF:
            return meta
        try:
            reader = _PdfReader(filepath)
            meta['pages'] = len(reader.pages)
            info = reader.metadata
            if info:
                if info.title:
                    meta['title'] = str(info.title).strip()
                if info.author:
                    meta['author'] = str(info.author).strip()
                if info.subject:
                    meta['subject'] = str(info.subject).strip()
                if info.creator:
                    meta['creator'] = str(info.creator).strip()
                if info.creation_date:
                    try:
                        meta['creation_date'] = info.creation_date.isoformat()
                    except Exception:
                        meta['creation_date'] = str(info.creation_date)
        except Exception:
            pass
        return meta

    # ── DOCX metadata ─────────────────────────────────────────────────────────
    @staticmethod
    def _extract_docx(filepath: str) -> dict:
        meta = {'_type': 'docx'}
        if not HAS_PYTHON_DOCX:
            return meta
        try:
            doc = _DocxDocument(filepath)
            props = doc.core_properties
            if props.title:
                meta['title'] = str(props.title).strip()
            if props.author:
                meta['author'] = str(props.author).strip()
            if props.subject:
                meta['subject'] = str(props.subject).strip()
            if props.last_modified_by:
                meta['last_modified_by'] = str(props.last_modified_by).strip()
            if props.keywords:
                meta['keywords'] = str(props.keywords).strip()
            if props.created:
                meta['created'] = props.created.isoformat()
            if props.modified:
                meta['modified'] = props.modified.isoformat()
            if props.revision:
                meta['revision'] = str(props.revision)
        except Exception:
            pass
        return meta

    # ── XLSX metadata ─────────────────────────────────────────────────────────
    @staticmethod
    def _extract_xlsx(filepath: str) -> dict:
        meta = {'_type': 'xlsx'}
        if not HAS_OPENPYXL:
            return meta
        try:
            wb = _load_workbook(filepath, read_only=True, data_only=True)
            props = wb.properties
            if props:
                if props.title:
                    meta['title'] = str(props.title).strip()
                if props.creator:
                    meta['author'] = str(props.creator).strip()
                if props.lastModifiedBy:
                    meta['last_modified_by'] = str(props.lastModifiedBy).strip()
                if props.subject:
                    meta['subject'] = str(props.subject).strip()
                if props.keywords:
                    meta['keywords'] = str(props.keywords).strip()
                if props.created:
                    meta['created'] = props.created.isoformat()
                if props.modified:
                    meta['modified'] = props.modified.isoformat()
            meta['sheet_count'] = len(wb.sheetnames)
            wb.close()
        except Exception:
            pass
        return meta

    # ── PPTX metadata (uses python-pptx if available, otherwise zipfile) ─────
    @staticmethod
    def _extract_pptx(filepath: str) -> dict:
        meta = {'_type': 'pptx'}
        try:
            import pptx as _pptx
            prs = _pptx.Presentation(filepath)
            props = prs.core_properties
            if props.title:
                meta['title'] = str(props.title).strip()
            if props.author:
                meta['author'] = str(props.author).strip()
            if props.subject:
                meta['subject'] = str(props.subject).strip()
            if props.last_modified_by:
                meta['last_modified_by'] = str(props.last_modified_by).strip()
            if props.created:
                meta['created'] = props.created.isoformat()
            if props.modified:
                meta['modified'] = props.modified.isoformat()
            meta['slide_count'] = len(prs.slides)
        except ImportError:
            # Fallback: extract slide count from ZIP contents
            try:
                import zipfile
                with zipfile.ZipFile(filepath, 'r') as zf:
                    slides = [n for n in zf.namelist()
                              if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
                    meta['slide_count'] = len(slides)
            except Exception:
                pass
        except Exception:
            pass
        return meta


# ══════════════════════════════════════════════════════════════════════════════
# ARCHIVE PEEK — Inspect archive contents for smart classification
# ══════════════════════════════════════════════════════════════════════════════


class ArchivePeeker:
    """Inspects ZIP/RAR/7z contents and classifies archives by what's inside."""

    @staticmethod
    def peek(filepath) -> dict:
        """Returns {'file_count': N, 'total_size': N, 'extensions': Counter, 'names': list[:20]}"""
        result = {'file_count': 0, 'total_size': 0, 'extensions': Counter(), 'names': []}
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == '.zip':
                with zipfile.ZipFile(filepath, 'r') as zf:
                    for info in zf.infolist():
                        if info.is_dir():
                            continue
                        result['file_count'] += 1
                        result['total_size'] += info.file_size
                        fext = os.path.splitext(info.filename)[1].lower()
                        if fext:
                            result['extensions'][fext] += 1
                        if len(result['names']) < 20:
                            result['names'].append(info.filename)
            elif ext == '.rar' and HAS_RARFILE:
                with _rarfile.RarFile(filepath, 'r') as rf:
                    for info in rf.infolist():
                        if info.is_dir():
                            continue
                        result['file_count'] += 1
                        result['total_size'] += info.file_size
                        fext = os.path.splitext(info.filename)[1].lower()
                        if fext:
                            result['extensions'][fext] += 1
                        if len(result['names']) < 20:
                            result['names'].append(info.filename)
            elif ext == '.7z' and HAS_PY7ZR:
                with _py7zr.SevenZipFile(filepath, 'r') as sz:
                    for name, bio in sz.read().items():
                        result['file_count'] += 1
                        result['total_size'] += bio.getbuffer().nbytes if hasattr(bio, 'getbuffer') else 0
                        fext = os.path.splitext(name)[1].lower()
                        if fext:
                            result['extensions'][fext] += 1
                        if len(result['names']) < 20:
                            result['names'].append(name)
        except Exception:
            pass
        return result

    @staticmethod
    def classify_contents(peek_result) -> tuple:
        """Returns (category, confidence) based on dominant extensions inside."""
        exts = peek_result.get('extensions', Counter())
        if not exts:
            return ('Archives', 50)
        top_ext = exts.most_common(1)[0][0] if exts else ''
        # Map dominant extension to category
        code_exts = {'.py', '.js', '.ts', '.java', '.c', '.cpp', '.go', '.rs', '.rb', '.php', '.html', '.css'}
        img_exts = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.psd', '.svg'}
        doc_exts = {'.pdf', '.doc', '.docx', '.txt', '.md', '.xlsx', '.csv'}
        audio_exts = {'.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a'}
        video_exts = {'.mp4', '.mkv', '.avi', '.mov', '.wmv', '.webm'}
        font_exts = {'.ttf', '.otf', '.woff', '.woff2'}
        total = sum(exts.values())
        for cat, ext_set in [('Code', code_exts), ('Images', img_exts), ('Documents', doc_exts),
                              ('Audio', audio_exts), ('Videos', video_exts), ('Fonts', font_exts)]:
            matched = sum(exts.get(e, 0) for e in ext_set)
            if matched > total * 0.5:
                return (cat, min(90, 60 + int(matched / total * 30)))
        return ('Archives', 55)


# ══════════════════════════════════════════════════════════════════════════════
# RULE ENGINE — User-defined classification rules
# ══════════════════════════════════════════════════════════════════════════════

_RULES_FILE = os.path.join(_APP_DATA_DIR, 'rules.json')


